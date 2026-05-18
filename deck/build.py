"""Build slide deck: From Rule-Taking to Rule-Making Chemistry."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# Palette
INK       = RGBColor(0x1A, 0x23, 0x32)   # deep ink/navy
CREAM     = RGBColor(0xFA, 0xF6, 0xF0)   # warm cream
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TERRA     = RGBColor(0xD9, 0x77, 0x57)   # terracotta — rule-makers
STEEL     = RGBColor(0x4F, 0x88, 0xB5)   # steel blue — rule-takers
SLATE     = RGBColor(0x6B, 0x72, 0x80)   # muted
GRAPH     = RGBColor(0x33, 0x3A, 0x44)   # body text
AMBER     = RGBColor(0xE8, 0xB0, 0x4A)
SOFTBG    = RGBColor(0xF4, 0xE9, 0xDC)   # soft warm card
LINE_TINT = RGBColor(0xE7, 0xDD, 0xCE)

# Layout
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

HEADER_FONT = "Georgia"
BODY_FONT   = "Calibri"

FIG_DIR = "../paper/figures"

prs = Presentation()
prs.slide_width  = Inches(SLIDE_W_IN)
prs.slide_height = Inches(SLIDE_H_IN)

BLANK = prs.slide_layouts[6]


def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, x, y, w, h, text, *,
             font=BODY_FONT, size=18, color=GRAPH,
             bold=False, italic=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, char_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
        if char_spacing is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set('spc', str(char_spacing))
    return tb


def add_rule(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_card(slide, x, y, w, h, fill=WHITE, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = border
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def page_number(slide, n, total, color=SLATE):
    add_text(slide, SLIDE_W_IN - 1.0, SLIDE_H_IN - 0.45, 0.8, 0.3,
             f"{n} / {total}", size=9, color=color, align=PP_ALIGN.RIGHT,
             font=BODY_FONT)


def header(slide, kicker, title, color_kicker=TERRA, color_title=INK):
    if kicker:
        add_text(slide, 0.7, 0.55, 12, 0.35, kicker.upper(),
                 size=11, color=color_kicker, bold=True, font=BODY_FONT,
                 char_spacing=300)
    add_text(slide, 0.7, 0.85, 12, 1.0, title,
             size=32, color=color_title, bold=True, font=HEADER_FONT)


TOTAL = 14
n = 0


# ─────────────────────────────────────────────────────────────────────
# Slide 1 — Title (dark)
# ─────────────────────────────────────────────────────────────────────
n += 1
s = prs.slides.add_slide(BLANK)
add_bg(s, INK)

# Subtle terracotta vertical mark
add_rule(s, 0.7, 2.8, 0.08, 1.8, TERRA)

add_text(s, 0.95, 2.7, 11, 0.45, "ORIGIN OF LIFE  ·  THEORY",
         size=12, color=TERRA, bold=True, font=BODY_FONT, char_spacing=400)

add_text(s, 0.95, 3.15, 11.5, 1.6,
         "From Rule-Taking",
         size=54, color=CREAM, bold=True, font=HEADER_FONT)
add_text(s, 0.95, 3.95, 11.5, 1.6,
         "to Rule-Making Chemistry",
         size=54, color=CREAM, bold=True, italic=True, font=HEADER_FONT)

add_text(s, 0.95, 5.15, 11.5, 0.5,
         "A minimal population–environment feedback model",
         size=18, color=LINE_TINT, italic=True, font=BODY_FONT)

# Footer author block
add_rule(s, 0.7, 6.55, 0.6, 0.03, TERRA)
add_text(s, 0.7, 6.65, 11.5, 0.35,
         "Celia Blanco", size=14, color=CREAM, bold=True, font=BODY_FONT)
add_text(s, 0.7, 6.95, 11.5, 0.35,
         "Centro de Astrobiología  ·  CSIC–INTA",
         size=11, color=SLATE, font=BODY_FONT)


# ─────────────────────────────────────────────────────────────────────
# Slide 2 — The opening question (cream)
# ─────────────────────────────────────────────────────────────────────
n += 1
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
header(s, "the question", "")
page_number(s, n, TOTAL)

# Centered hero quote
add_text(s, 1.2, 2.5, 11, 1.4,
         "When does chemistry stop",
         size=44, color=INK, bold=True, font=HEADER_FONT, align=PP_ALIGN.CENTER)
add_text(s, 1.2, 3.3, 11, 1.4,
         "adapting to its conditions —",
         size=44, color=INK, bold=True, font=HEADER_FONT, align=PP_ALIGN.CENTER)
add_text(s, 1.2, 4.1, 11, 1.4,
         "and start authoring them?",
         size=44, color=TERRA, bold=True, italic=True,
         font=HEADER_FONT, align=PP_ALIGN.CENTER)

add_rule(s, SLIDE_W_IN/2 - 0.4, 5.6, 0.8, 0.03, TERRA)


# ─────────────────────────────────────────────────────────────────────
# Slide 3 — What classical models leave fixed
# ─────────────────────────────────────────────────────────────────────
n += 1
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
header(s, "context", "Classical origin-of-life models")
page_number(s, n, TOTAL)

# two-column
col_y = 2.4
col_h = 4.2
gap_x = 0.4
col_w = (SLIDE_W_IN - 1.4 - gap_x) / 2

# Left card — what's modeled
add_card(s, 0.7, col_y, col_w, col_h, fill=WHITE)
add_rule(s, 0.7, col_y, 0.08, col_h, STEEL)
add_text(s, 0.95, col_y + 0.35, col_w - 0.35, 0.4,
         "WELL STUDIED", size=11, color=STEEL, bold=True,
         font=BODY_FONT, char_spacing=300)
add_text(s, 0.95, col_y + 0.75, col_w - 0.35, 0.6,
         "Replicators", size=28, color=INK, bold=True, font=HEADER_FONT)
add_text(s, 0.95, col_y + 1.6, col_w - 0.35, 2.4,
         ["Autocatalytic sets",
          "Hypercycles",
          "Protocells"],
         size=18, color=GRAPH, font=BODY_FONT)
add_text(s, 0.95, col_y + 3.45, col_w - 0.35, 0.5,
         "Fidelity, replication, inheritance.",
         size=13, color=SLATE, italic=True, font=BODY_FONT)

# Right card — what's left fixed
rx = 0.7 + col_w + gap_x
add_card(s, rx, col_y, col_w, col_h, fill=WHITE)
add_rule(s, rx, col_y, 0.08, col_h, TERRA)
add_text(s, rx + 0.25, col_y + 0.35, col_w - 0.35, 0.4,
         "USUALLY FIXED", size=11, color=TERRA, bold=True,
         font=BODY_FONT, char_spacing=300)
add_text(s, rx + 0.25, col_y + 0.75, col_w - 0.35, 0.6,
         "The environment", size=28, color=INK, bold=True, font=HEADER_FONT)
add_text(s, rx + 0.25, col_y + 1.6, col_w - 0.35, 2.4,
         ["Reaction conditions imposed externally",
          "Not shaped by the chemistry it hosts",
          "No feedback into selection"],
         size=18, color=GRAPH, font=BODY_FONT)
add_text(s, rx + 0.25, col_y + 3.45, col_w - 0.35, 0.5,
         "But life is defined by reshaping its surroundings.",
         size=13, color=SLATE, italic=True, font=BODY_FONT)


# ─────────────────────────────────────────────────────────────────────
# Slide 4 — Two chemistries (concept cartoon)
# ─────────────────────────────────────────────────────────────────────
n += 1
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
header(s, "the distinction", "Two chemistries")
page_number(s, n, TOTAL)

# Image: keep aspect ratio — original is roughly 2:1
img_w = 11.0
img_h = img_w * 0.42  # rough aspect
img_x = (SLIDE_W_IN - img_w) / 2
img_y = 2.1
s.shapes.add_picture(f"{FIG_DIR}/concept_cartoon.png",
                     Inches(img_x), Inches(img_y),
                     width=Inches(img_w))

# Caption below
cap_y = 6.3
add_text(s, 1.2, cap_y, 11, 0.5,
         "Rule-takers (R) depend on the environment.  "
         "Rule-makers (M) reshape it — at a cost.",
         size=16, color=GRAPH, italic=True, font=BODY_FONT,
         align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────
# Slide 5 — Minimal model schematic
# ─────────────────────────────────────────────────────────────────────
n += 1
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
header(s, "the model", "Minimal coupling: population ↔ environment")
page_number(s, n, TOTAL)

import math

def conn(slide, x1, y1, x2, y2, color, *, width=2.5, dashed=False):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1),
                                       Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    lnL = line.line._get_or_add_ln()
    if dashed:
        dash = etree.SubElement(lnL, qn('a:prstDash'))
        dash.set('val', 'dash')
    tail = etree.SubElement(lnL, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('h', 'med')
    return line

def circle(slide, cx_in, cy_in, r_in, fill, label, *, label_color=WHITE,
           label_size=30):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(cx_in - r_in), Inches(cy_in - r_in),
                                Inches(2*r_in), Inches(2*r_in))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    add_text(slide, cx_in - r_in, cy_in - r_in, 2*r_in, 2*r_in,
             label, size=label_size, color=label_color, bold=True,
             font=HEADER_FONT, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

# Triangle: E at the apex, R and M at the base.
# E reaches DOWN to both R and M.  Only M reaches BACK UP to E.
E_cx, E_cy, E_r = 6.667, 3.05, 0.95
R_cx, R_cy, R_r = 2.5, 5.55, 0.82
M_cx, M_cy, M_r = 10.83, 5.55, 0.82

# Draw E with "E" letter + small "environment" subtitle inside the circle
def circle_with_subtitle(slide, cx_in, cy_in, r_in, fill,
                         label, sub, *, label_size=40, sub_size=11):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(cx_in - r_in), Inches(cy_in - r_in),
                                Inches(2*r_in), Inches(2*r_in))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    add_text(slide, cx_in - r_in, cy_in - r_in - 0.10, 2*r_in, r_in + 0.30,
             label, size=label_size, color=WHITE, bold=True,
             font=HEADER_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
    add_text(slide, cx_in - r_in, cy_in + 0.05, 2*r_in, 0.32,
             sub, size=sub_size, color=WHITE, italic=True,
             font=BODY_FONT, align=PP_ALIGN.CENTER)

circle_with_subtitle(s, E_cx, E_cy, E_r, TERRA, "E", "environment",
                     label_size=40, sub_size=11)
circle(s, R_cx, R_cy, R_r, STEEL, "R", label_size=36)
circle(s, M_cx, M_cy, M_r, TERRA, "M", label_size=36)

# Labels beneath R and M (outside circles, below triangle base)
def node_label(slide, x_center, y, title, sub):
    add_text(slide, x_center - 1.8, y, 3.6, 0.32,
             title, size=12, color=INK, bold=True,
             font=BODY_FONT, align=PP_ALIGN.CENTER)
    add_text(slide, x_center - 2.0, y + 0.32, 4.0, 0.3,
             sub, size=10, color=SLATE, italic=True,
             font=BODY_FONT, align=PP_ALIGN.CENTER)

node_label(s, R_cx, R_cy + R_r + 0.10, "rule-taker",
           "fraction 1 – y")
node_label(s, M_cx, M_cy + M_r + 0.10, "rule-maker",
           "fraction y  ·  pays cost c")

# ----- Arrows -----
def edge(cx, cy, r, tx, ty):
    dx, dy = tx - cx, ty - cy
    L = math.hypot(dx, dy)
    return cx + r*dx/L, cy + r*dy/L

# Parallel-arrow pair between M and E
ux, uy = (E_cx - M_cx), (E_cy - M_cy)
L = math.hypot(ux, uy); ux, uy = ux/L, uy/L      # unit M→E
px, py = -uy, ux                                  # perpendicular (CCW from M→E)
offset = 0.32

M_edge = edge(M_cx, M_cy, M_r, E_cx, E_cy)
E_edge_M = edge(E_cx, E_cy, E_r, M_cx, M_cy)

# M → E (production, terracotta), shifted in +perp direction
conn(s,
     M_edge[0]   + offset*px, M_edge[1]   + offset*py,
     E_edge_M[0] + offset*px, E_edge_M[1] + offset*py,
     TERRA)
# E → M (feedback, steel), shifted in -perp direction
conn(s,
     E_edge_M[0] - offset*px, E_edge_M[1] - offset*py,
     M_edge[0]   - offset*px, M_edge[1]   - offset*py,
     STEEL)

# Labels for the two M↔E arrows — placed OUTSIDE each arrow (away from the pair)
mid_x_base = (M_edge[0] + E_edge_M[0]) / 2
mid_y_base = (M_edge[1] + E_edge_M[1]) / 2

# Label for M → E (TERRA): further in +perp from the terracotta arrow
lblA_x = mid_x_base + 2.6*offset*px
lblA_y = mid_y_base + 2.6*offset*py
add_text(s, lblA_x - 1.4, lblA_y - 0.18, 2.8, 0.36,
         "produces  (α)", size=13, color=TERRA, bold=True,
         font=BODY_FONT, align=PP_ALIGN.CENTER)

# Label for E → M (STEEL): further in -perp
lblB_x = mid_x_base - 2.6*offset*px
lblB_y = mid_y_base - 2.6*offset*py
add_text(s, lblB_x - 1.4, lblB_y - 0.18, 2.8, 0.36,
         "enhances  (p)", size=13, color=STEEL, bold=True,
         font=BODY_FONT, align=PP_ALIGN.CENTER)

# E → R (single dashed arrow — environment reaches R but R doesn't reach back)
E_edge_R = edge(E_cx, E_cy, E_r, R_cx, R_cy)
R_edge_E = edge(R_cx, R_cy, R_r, E_cx, E_cy)
conn(s, E_edge_R[0], E_edge_R[1], R_edge_E[0], R_edge_E[1],
     SLATE, width=1.6, dashed=True)

# Label for E → R — perpendicular-offset OUTWARD (away from the triangle's interior)
ER_ux, ER_uy = (R_cx - E_cx), (R_cy - E_cy)
LER = math.hypot(ER_ux, ER_uy); ER_ux, ER_uy = ER_ux/LER, ER_uy/LER
# perpendicular toward upper-right of the arrow (i.e. AWAY from R-M axis)
ER_px, ER_py = ER_uy, -ER_ux   # rotated 90° CW so it points "up-right" of E→R
mid_ER_x = (E_edge_R[0] + R_edge_E[0]) / 2
mid_ER_y = (E_edge_R[1] + R_edge_E[1]) / 2
add_text(s, mid_ER_x - 0.30*ER_px - 1.6,
            mid_ER_y - 0.30*ER_py - 0.18,
            3.2, 0.36,
         "no effect on R’s rate", size=12, color=SLATE, italic=True,
         font=BODY_FONT, align=PP_ALIGN.CENTER)

# Footnote
add_text(s, 0.7, 6.95, 12, 0.3,
         "R sits in E but ignores it (π_R = 1).   Only M’s replication rate responds to E.",
         size=13, color=SLATE, italic=True, font=BODY_FONT,
         align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────
# Slide 6 — Equations
# ─────────────────────────────────────────────────────────────────────
n += 1
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
header(s, "the dynamics", "Two coupled equations")
page_number(s, n, TOTAL)

# Equation card
eq_y = 2.6
add_card(s, 1.8, eq_y, 9.7, 3.0, fill=WHITE)
add_rule(s, 1.8, eq_y, 0.08, 3.0, INK)

# Eq 1
add_text(s, 2.4, eq_y + 0.45, 9.0, 0.5,
         "POPULATION    (replicator equation)",
         size=11, color=SLATE, bold=True, font=BODY_FONT, char_spacing=300)
add_text(s, 2.4, eq_y + 0.85, 9.0, 0.8,
         "ẏ  =  y (1 – y) · ( s – c + p E – 1 )",
         size=26, color=INK, bold=True, italic=True, font="Georgia")

# Eq 2
add_text(s, 2.4, eq_y + 1.85, 9.0, 0.5,
         "ENVIRONMENT    (production – decay)",
         size=11, color=SLATE, bold=True, font=BODY_FONT, char_spacing=300)
add_text(s, 2.4, eq_y + 2.25, 9.0, 0.8,
         "Ė  =  α y  –  β E",
         size=26, color=INK, bold=True, italic=True, font="Georgia")

# Subtitle
add_text(s, 0.7, 6.0, 12, 0.5,
         "M produces E.  A larger E raises M’s replication rate — more M, more E.",
         size=15, color=GRAPH, italic=True, font=BODY_FONT,
         align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────
# Slide 7 — One number controls everything
# ─────────────────────────────────────────────────────────────────────
n += 1
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
header(s, "the threshold", "One number controls the regime")
page_number(s, n, TOTAL)

# Big inequality, centered — built with subscript runs
tb = s.shapes.add_textbox(Inches(1.0), Inches(2.7), Inches(11.3), Inches(1.0))
tf = tb.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = 0
tf.margin_top = tf.margin_bottom = 0
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER

def eqrun(p, text, *, size=34, italic=True, bold=True, color=INK, sub=False):
    r = p.add_run(); r.text = text
    f = r.font
    f.name = "Georgia"
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color
    if sub:
        rPr = r._r.get_or_add_rPr()
        rPr.set('baseline', '-25000')
    return r

eqrun(p, "p")
eqrun(p, "eff", sub=True)
eqrun(p, "  =  p α / β       >       p")
eqrun(p, "c", sub=True)
eqrun(p, "  =  c + 1 – s")

# Two captions below
cap_y = 4.3
cap_w = 5.3
add_text(s, 1.3, cap_y, cap_w, 0.5,
         "EFFECTIVE FEEDBACK",
         size=11, color=TERRA, bold=True, font=BODY_FONT,
         char_spacing=300, align=PP_ALIGN.CENTER)
add_text(s, 1.3, cap_y + 0.45, cap_w, 0.7,
         "how strongly modifiers help,",
         size=15, color=GRAPH, font=BODY_FONT, align=PP_ALIGN.CENTER)
add_text(s, 1.3, cap_y + 0.75, cap_w, 0.7,
         "weighted by how long they last",
         size=15, color=GRAPH, font=BODY_FONT, align=PP_ALIGN.CENTER)

add_text(s, SLIDE_W_IN - 1.3 - cap_w, cap_y, cap_w, 0.5,
         "COST BARRIER",
         size=11, color=STEEL, bold=True, font=BODY_FONT,
         char_spacing=300, align=PP_ALIGN.CENTER)
add_text(s, SLIDE_W_IN - 1.3 - cap_w, cap_y + 0.45, cap_w, 0.7,
         "what rule-makers give up",
         size=15, color=GRAPH, font=BODY_FONT, align=PP_ALIGN.CENTER)
add_text(s, SLIDE_W_IN - 1.3 - cap_w, cap_y + 0.75, cap_w, 0.7,
         "to modify the environment",
         size=15, color=GRAPH, font=BODY_FONT, align=PP_ALIGN.CENTER)

# Punchline
add_rule(s, SLIDE_W_IN/2 - 0.4, 6.3, 0.8, 0.03, TERRA)
add_text(s, 1.0, 6.45, 11.3, 0.5,
         "Cross it, and a second stable state appears.",
         size=18, color=INK, italic=True, bold=True, font=HEADER_FONT,
         align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────
# Slide 8 — Phase diagram
# ─────────────────────────────────────────────────────────────────────
n += 1
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
header(s, "phase diagram", "Cost versus feedback strength")
page_number(s, n, TOTAL)

# image on the right; takeaways on the left
img_w = 6.4
img_h = img_w * (1500/2200)  # original ~2200x1500
img_x = SLIDE_W_IN - img_w - 0.7
img_y = 2.1
s.shapes.add_picture(f"{FIG_DIR}/regimes.png",
                     Inches(img_x), Inches(img_y),
                     width=Inches(img_w))

# Left text
tx = 0.7
ty = 2.3
add_text(s, tx, ty, 5.2, 0.5,
         "ABOVE THE LINE", size=11, color=TERRA, bold=True,
         font=BODY_FONT, char_spacing=300)
add_text(s, tx, ty + 0.45, 5.2, 1.4,
         "Rule-takers and rule-makers each form a stable state.",
         size=18, color=INK, bold=True, font=HEADER_FONT)

add_text(s, tx, ty + 1.95, 5.2, 0.5,
         "BELOW THE LINE", size=11, color=STEEL, bold=True,
         font=BODY_FONT, char_spacing=300)
add_text(s, tx, ty + 2.4, 5.2, 1.4,
         "Only rule-taking persists — modifiers can’t pay for themselves.",
         size=18, color=INK, bold=True, font=HEADER_FONT)

add_text(s, tx, 6.3, 5.2, 0.5,
         "Stronger feedback compensates for higher cost — linearly.",
         size=13, color=SLATE, italic=True, font=BODY_FONT)


# ─────────────────────────────────────────────────────────────────────
# Slide 9 — Nullclines / geometry of bistability
# ─────────────────────────────────────────────────────────────────────
n += 1
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
header(s, "geometry", "A saddle divides the basin")
page_number(s, n, TOTAL)

img_w = 6.0
img_x = SLIDE_W_IN - img_w - 0.7
img_y = 2.0
s.shapes.add_picture(f"{FIG_DIR}/2D_nullcline.png",
                     Inches(img_x), Inches(img_y),
                     width=Inches(img_w))

# Left text
tx = 0.7
add_text(s, tx, 2.2, 5.4, 0.5,
         "TWO ATTRACTORS",
         size=11, color=TERRA, bold=True, font=BODY_FONT, char_spacing=300)
add_text(s, tx, 2.6, 5.4, 1.5,
         "A₁ = rule-takers only.   A₂ = rule-makers dominate.",
         size=18, color=INK, bold=True, font=HEADER_FONT)

add_text(s, tx, 4.4, 5.4, 0.5,
         "ONE SADDLE",
         size=11, color=STEEL, bold=True, font=BODY_FONT, char_spacing=300)
add_text(s, tx, 4.8, 5.4, 1.5,
         "Its stable manifold draws the separatrix between the two fates.",
         size=18, color=INK, bold=True, font=HEADER_FONT)


# ─────────────────────────────────────────────────────────────────────
# Slide 10 — Without environmental memory
# ─────────────────────────────────────────────────────────────────────
n += 1
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
header(s, "case 1 — no memory", "Composition alone decides")
page_number(s, n, TOTAL)

# Image full width (aspect ~0.438)
img_w = 9.4
img_x = (SLIDE_W_IN - img_w) / 2
img_y = 2.15
s.shapes.add_picture(f"{FIG_DIR}/timeseries_fan.png",
                     Inches(img_x), Inches(img_y),
                     width=Inches(img_w))

# Bottom caption
add_text(s, 1.0, 6.55, 11.3, 0.5,
         "When the environment relaxes instantly, only a sharp frequency threshold matters.",
         size=15, color=GRAPH, italic=True, font=BODY_FONT,
         align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────
# Slide 11 — With environmental memory
# ─────────────────────────────────────────────────────────────────────
n += 1
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
header(s, "case 2 — environmental memory", "History reshapes the basin")
page_number(s, n, TOTAL)

img_w = 10.6
img_x = (SLIDE_W_IN - img_w) / 2
img_y = 2.15
s.shapes.add_picture(f"{FIG_DIR}/extended_fan.png",
                     Inches(img_x), Inches(img_y),
                     width=Inches(img_w))

add_text(s, 1.0, 6.55, 11.3, 0.5,
         "Because E remembers, a favourable past lets even rare rule-makers invade.",
         size=15, color=GRAPH, italic=True, font=BODY_FONT,
         align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────
# Slide 12 — Chemical realization
# ─────────────────────────────────────────────────────────────────────
n += 1
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
header(s, "chemical realization", "Same dynamics, real molecules")
page_number(s, n, TOTAL)

# Left: reaction list in a card
card_x = 0.7
card_y = 2.05
card_w = 4.9
card_h = 4.6
add_card(s, card_x, card_y, card_w, card_h, fill=WHITE)
add_rule(s, card_x, card_y, 0.08, card_h, TERRA)

add_text(s, card_x + 0.25, card_y + 0.3, card_w - 0.35, 0.4,
         "MINIMAL REACTION NETWORK",
         size=10, color=TERRA, bold=True, font=BODY_FONT, char_spacing=300)

rxns = [
    "R + S  →  2R",
    "M + S  →  2M",
    "M + S + B  →  2M + B",
    "M + S  →  M + B",
    "B  →  ∅",
]
labels = [
    "rule-taker replication",
    "rule-maker replication",
    "modifier-assisted (catalytic)",
    "modifier production (costly)",
    "modifier decay",
]
ry = card_y + 0.85
for rxn, lbl in zip(rxns, labels):
    add_text(s, card_x + 0.3, ry, card_w - 0.5, 0.36,
             rxn, size=14, color=INK, bold=True, font="Georgia")
    add_text(s, card_x + 0.3, ry + 0.36, card_w - 0.5, 0.32,
             lbl, size=10, color=SLATE, italic=True, font=BODY_FONT)
    ry += 0.72

# Right: chemical_feedback figure (aspect 0.40)
img_w = 7.2
img_h = img_w * 0.40
img_x = SLIDE_W_IN - img_w - 0.4
img_y = card_y + (card_h - img_h) / 2
s.shapes.add_picture(f"{FIG_DIR}/chemical_feedback_two_panel.png",
                     Inches(img_x), Inches(img_y),
                     width=Inches(img_w))


# ─────────────────────────────────────────────────────────────────────
# Slide 13 — Take-home (three points)
# ─────────────────────────────────────────────────────────────────────
n += 1
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
header(s, "take home", "Three ideas to keep")
page_number(s, n, TOTAL)

# Three cards in a row
items = [
    ("01", "Persistence creates inheritance",
     "Modifiers that outlast a generation transmit conditions, not genes — a chemistry-level form of heredity.",
     TERRA),
    ("02", "History opens basins",
     "A favourable environment lowers the threshold for rule-makers — even when they start rare.",
     STEEL),
    ("03", "Autonomy is conditional",
     "Bistability is not guaranteed: feedback must be strong, persistent, and seeded by the right history.",
     INK),
]
gap = 0.35
card_w = (SLIDE_W_IN - 1.4 - 2*gap) / 3
card_h = 4.4
cy = 2.3
cx = 0.7
for num, ttl, body, accent in items:
    add_card(s, cx, cy, card_w, card_h, fill=WHITE)
    add_rule(s, cx, cy, 0.08, card_h, accent)
    add_text(s, cx + 0.3, cy + 0.4, card_w - 0.5, 0.7,
             num, size=36, color=accent, bold=True, font=HEADER_FONT, italic=True)
    add_text(s, cx + 0.3, cy + 1.3, card_w - 0.5, 1.4,
             ttl, size=20, color=INK, bold=True, font=HEADER_FONT)
    add_text(s, cx + 0.3, cy + 2.7, card_w - 0.5, 1.5,
             body, size=13, color=GRAPH, font=BODY_FONT)
    cx += card_w + gap


# ─────────────────────────────────────────────────────────────────────
# Slide 14 — Closing (dark)
# ─────────────────────────────────────────────────────────────────────
n += 1
s = prs.slides.add_slide(BLANK)
add_bg(s, INK)

add_rule(s, 0.7, 2.8, 0.08, 1.8, TERRA)

add_text(s, 0.95, 2.7, 11.5, 0.5,
         "CLOSING THOUGHT",
         size=11, color=TERRA, bold=True, font=BODY_FONT, char_spacing=400)

add_text(s, 0.95, 3.1, 11.5, 1.0,
         "Replication makes a chemistry persist.",
         size=30, color=CREAM, font=HEADER_FONT)
add_text(s, 0.95, 3.75, 11.5, 1.0,
         "Feedback lets it set the conditions for its own persistence.",
         size=30, color=CREAM, italic=True, font=HEADER_FONT)

add_text(s, 0.95, 5.0, 11.5, 0.5,
         "That second step is rare — and may be where chemistry began to look like life.",
         size=16, color=LINE_TINT, italic=True, font=BODY_FONT)

add_rule(s, 0.7, 6.55, 0.6, 0.03, TERRA)
add_text(s, 0.7, 6.65, 11.5, 0.35,
         "Thank you.", size=14, color=CREAM, bold=True, font=BODY_FONT)
add_text(s, 0.7, 6.95, 11.5, 0.35,
         "celia.blanco@cab.inta-csic.es",
         size=11, color=SLATE, font=BODY_FONT)


# Save
out = "RuleMaking_Chemistry.pptx"
prs.save(out)
print(f"Wrote {out}  ({n} slides)")
